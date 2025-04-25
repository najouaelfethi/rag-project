from langchain_community.document_loaders import PyPDFLoader,UnstructuredWordDocumentLoader,UnstructuredPowerPointLoader
from langchain.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter,MarkdownHeaderTextSplitter,SentenceTransformersTokenTextSplitter
from langchain_experimental.text_splitter import SemanticChunker
import os
from langchain_openai import ChatOpenAI
from dotenv import load_dotenv
from langchain_core.messages import SystemMessage, HumanMessage
from langchain_openai import OpenAIEmbeddings
from pptx import Presentation


load_dotenv()

openai_key = os.getenv("OPENAI_API_KEY")
openai_org = os.getenv("OPENAI_ORG_ID")

model = OpenAIEmbeddings(openai_api_key=openai_key, openai_organization=openai_org)

#---------------CHUNKING-----------------------

def chunk_pdf(file_path): #Splitting using Semantic Chunker => split by meaning
    docs = PyPDFLoader(file_path).load()#load and variable stores content of PDF
    splitter = SemanticChunker(model)#split doc into chunks, each chunk is of size 600 and each new chunk repeats 100 characters fro the previous one(chunk)=>AI will remeber context better
    chunks = splitter.split_documents(docs)
    return chunks
     

def chunk_docx(file_path):#Hybrid splitting: semantic(using mMarkDown...) & character splitting if no heading found
    docs = UnstructuredWordDocumentLoader(file_path).load()#load and variable stores content of doc
    raw_text = "\n".join([doc.page_content for doc in docs])
    heading_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=[("#","title"),("##","section")])#splitting by headings(titles) and label chunk as section=> semantic splitting it split by meaningful document structure not just size
    heading_chunks = heading_splitter.split_text(raw_text)
    if len(heading_chunks)<=1:#means we have probably any clear headings
        print("no heading detected")
        char_splitter = RecursiveCharacterTextSplitter(chunk_size=600, chunk_overlap=100)
        chunks = char_splitter.split_documents(docs)
    else:
        chunks=heading_chunks
    return chunks


def chunk_txt(file_path):#Splitting using Semantic Chunker => split by meaning
    with open(file_path, "r") as f:
        raw_text = f.read()
        docs = [Document(page_content=raw_text)]
        splitter = SemanticChunker(model)
        chunks = splitter.split_documents(docs)
    return chunks


def chunk_pptx(file_path):#splitting using pptx library, 1 slide=1 chunk, ppt stores slides as seperate docs
    prs = Presentation(file_path)
    slide_texts=[]#storing text of each slide
    for i, slide in enumerate(prs.slides, start=1):
        text=""
        for shape in slide.shapes:
            if hasattr(shape,"text"):#checking if an object has specefic feature(text in this case)
                text+=shape.text + "\n"
        slide_texts.append(text.strip())
    return slide_texts

#------------TESTING----------------------

def test_pdf():
    file_name = "./documents/attention_pdf.pdf"
    chunks = chunk_pdf(file_name)
    for i, doc in enumerate(chunks[:10]):#limit 10 chunks
        print(f"chunk {i+1} content:\n{doc.page_content}")
        #feedback = evaluate_chunks(doc.page_contefor i, doc in enumerate(chunks)
        #print(f"\nGPT Feedback:\n{feedback}")

def test_txt():
    file_name = "./documents/machine learning notes.txt"
    chunks = chunk_txt(file_name)
    for i, doc in enumerate(chunks[:10]):#limit 10 chunks
        print(f"chunk {i+1} content:\n{doc.page_content}")
        #feedback = evaluate_chunks(doc.page_contefor i, doc in enumerate(chunks)
        #print(f"\nGPT Feedback:\n{feedback}")

def test_docx():
    file_name = "./documents/paper Optimization.docx"
    chunks = chunk_docx(file_name)
    (file_name)
    for i, doc in enumerate(chunks[:10]):
        print(f"\n Chunk {i+1} content:")
        print(doc.page_content)
        #feedback = evaluate_chunks(doc.page_content)
        #print(f"\nGPT Feedback:\n{feedback}")

def test_pptx():
    file_name = "./documents/MachineLearning FINAL.pptx"
    chunks = chunk_pptx(file_name)
    for i, chunk in enumerate(chunks,1):
        print(f"\n Slide {i} content:\n{chunk}")


#--------Evaluating chunks with GPT4---------------

def evaluate_chunks(chunk_text):
    prompt=f""" 
you are an design experct, rate the folowwing text chunk from 1(bad) to 5(excellent) without any extra informations

Chunk:
{chunk_text}
"""
    llm = ChatOpenAI(temperature=0, model_name='gpt-4o-mini',openai_api_key=openai_key, openai_organization=openai_org)#temperature controls creativity of the model
    response = llm.invoke([
          SystemMessage(content="You are senior e-learning content reviewer and educational designer"),
          HumanMessage(content=prompt)])
    
    return response.content.strip()


if __name__=='__main__':
    test_pptx()



    