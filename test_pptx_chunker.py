from pptx import Presentation

def main():
    file_path = "./documents/MachineLearning FINAL.pptx"
    prs = Presentation(file_path)
    slide_chunks = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []

        # Titre
        if slide.shapes.title and slide.shapes.title.text:
            slide_text.append(f"# {slide.shapes.title.text.strip()}")

        # Autres éléments textuels
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        # Notes de la slide
        notes = ""
        if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide and hasattr(notes_slide, "notes_text_frame") and notes_slide.notes_text_frame:
                notes_frame = notes_slide.notes_text_frame
                if notes_frame.text:
                    notes = notes_frame.text.strip()
                    slide_text.append(f"NOTES: {notes}")

        # Regrouper le contenu
        content = "\n\n".join(slide_text)
        if content:
            slide_chunks.append({
                "slide_number": i,
                "content": content,
                "notes": notes
            })

    # Affichage des résultats
    for chunk in slide_chunks:
        print(f"\n--- Slide {chunk['slide_number']} ---")
        print(chunk["content"])

if __name__ == "__main__":
    main()
